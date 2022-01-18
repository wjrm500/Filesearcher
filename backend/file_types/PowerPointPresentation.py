from pptx import Presentation

from file_types.File import File

class PowerPointPresentation(File):
    extension = 'pptx'

    def get_text(self) -> str:
        presentation = Presentation(self.filename)
        text_elements = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame and (text := shape.text) != '':
                    text_elements.append(text)
                if shape.has_table:
                    table = shape.table
                    for row in range(0, len(table.rows)):
                        for column in range(0, len(table.columns)):
                            cell = table.cell(row, column)
                            text = ' '.join([run.text for paragraph in cell.text_frame.paragraphs for run in paragraph.runs])
                            if text != '':
                                text_elements.append(text)
        return '\n'.join(text_elements)